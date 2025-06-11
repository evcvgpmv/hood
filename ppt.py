from pptx import Presentation

# Create a new presentation
prs = Presentation()

# Function to add a slide
def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "\n".join(content_lines)

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Ask Your On-Premise Data"
slide.placeholders[1].text = (
    "Harnessing GPT-4o, Chainlit, and Azure AI for Secure, Conversational Data Access\n"
    "Presented by: [Your Name]  \nDate: [Insert Date]"
)

# Add content slides
slides_content = [
    ("Problem Statement", [
        "- Business data is often siloed in on-premise systems.",
        "- Non-technical users struggle to access and query structured data efficiently.",
        "- Manual SQL writing is time-consuming and error-prone.",
        "",
        "🔍 Goal: Enable users to ask questions in natural language and get accurate answers from on-premise data, securely."
    ]),
    ("Solution Overview", [
        "🧠 Conversational AI for On-Prem Data",
        "A lightweight, secure, and intelligent assistant that allows users to:",
        "- Ask natural language questions.",
        "- Retrieve insights from on-premise databases.",
        "- Get human-like, accurate, and contextual responses.",
        "",
        "Powered by: 🔗 Chainlit, 🧠 GPT-4o, ⚙️ Azure AI Foundry, 🧩 Function Calling, 🗄️ SQLite Database"
    ]),
    ("Architecture Diagram", [
        "📊 System Architecture:",
        "- Frontend: Chainlit (Chat-based UI)",
        "- LLM: GPT-4o via Azure AI Foundry",
        "- Function Tool: Custom tool calling for SQL generation",
        "- Database: On-premise SQLite instance",
        "- Security: Local execution, no sensitive data leaves premises",
        "",
        "[Insert visual diagram here]"
    ]),
    ("Tech Stack Breakdown", [
        "| Component | Tool Used | Purpose |",
        "|-----------|-----------|---------|",
        "| Chat UI   | Chainlit  | User interaction interface |",
        "| LLM       | GPT-4o via Azure AI Foundry | Natural language processing |",
        "| Logic     | Function Tooling | Convert queries → SQL |",
        "| Storage   | SQLite    | Lightweight, on-prem database |",
        "| Hosting   | Local / VM | On-premise secure environment |"
    ]),
    ("Function Tool in Action", [
        "⚙️ Function Tool Highlights:",
        "- Dynamically maps user intent to SQL queries.",
        "- Ensures validated, schema-aware query generation.",
        "- Reduces hallucination risk with structured function parameters.",
        "",
        "💡 Example:",
        'User: "Show me the top 5 products by revenue last month"',
        "→ Generates safe SQL based on schema context."
    ]),
    ("Chainlit UI Demo", [
        "🖼️ Interactive Chat Interface:",
        "- Easy-to-use UI for business users",
        "- Real-time responses",
        "- SQL output shown for transparency (optional)",
        "",
        "[Add screenshots or link to a demo video here]"
    ]),
    ("Key Benefits", [
        "✅ Business Impact:",
        "- Democratizes access to data",
        "- Boosts productivity of business users",
        "- Reduces load on data/BI teams",
        "- Highly secure (on-prem deployment)"
    ]),
    ("Challenges & Mitigation", [
        "🚧 Challenges Addressed:",
        "| Challenge            | Solution                                  |",
        "|----------------------|-------------------------------------------|",
        "| LLM hallucination    | Function calling + schema-aware prompts   |",
        "| Data security        | Entire stack runs on-prem                 |",
        "| Usability            | Intuitive UI with Chainlit                |",
        "| Performance          | Optimized SQLite queries                  |"
    ]),
    ("Roadmap & Next Steps", [
        "🛣️ Roadmap:",
        "- ✅ PoC completed",
        "- 🔄 Add support for multiple databases (PostgreSQL, MS SQL)",
        "- 🔐 Role-based access to sensitive queries",
        "- 📈 Log & analyze user query patterns"
    ]),
    ("Q&A + Call to Action", [
        "🙋 Questions?",
        "",
        "👉 Let’s explore how this can be scaled across departments or integrated into existing analytics platforms."
    ])
]

for title, content in slides_content:
    add_slide(title, content)

# Save the presentation
prs.save("Ask_Your_OnPrem_Data_Presentation.pptx")
print("Presentation saved as 'Ask_Your_OnPrem_Data_Presentation.pptx'")
